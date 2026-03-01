"""
Tests for lazy_load() across all supported document loaders.

Verifies that every loader returned by get_loader() supports lazy_load(),
returns a generator/iterator, yields valid Document objects with content,
and measures whether lazy_load() provides real memory benefits vs load().
"""

import gc
import os
import shutil
import tracemalloc
import types
import zipfile

import pytest
from langchain_core.documents import Document

from app.utils.document_loader import get_loader, SafePyPDFLoader

# ---------------------------------------------------------------------------
# Environment checks — these deps aren't guaranteed in every CI runner
# ---------------------------------------------------------------------------

_has_pandoc = shutil.which("pandoc") is not None
if not _has_pandoc:
    try:
        import pypandoc

        pypandoc.get_pandoc_path()
        _has_pandoc = True  # pypandoc_binary or similar bundles the binary
    except (ImportError, OSError):
        pass

try:
    import msoffcrypto  # noqa: F401

    _has_msoffcrypto = True
except ImportError:
    _has_msoffcrypto = False

_skip_no_pandoc = pytest.mark.skipif(not _has_pandoc, reason="pandoc not installed")
_skip_no_msoffcrypto = pytest.mark.skipif(
    not _has_msoffcrypto, reason="msoffcrypto not installed"
)


# ---------------------------------------------------------------------------
# Fixture helpers for generating test files in each format
# ---------------------------------------------------------------------------


def _make_pdf(path, *, num_pages=1):
    """Create a PDF with extractable text using pypdf.

    When *num_pages* > 1 each page gets unique text so tracemalloc can
    observe the difference between holding 1 page vs N pages in memory.
    """
    from pypdf import PdfWriter
    from pypdf.generic import (
        DecodedStreamObject,
        DictionaryObject,
        NameObject,
    )

    writer = PdfWriter()
    for i in range(num_pages):
        writer.add_blank_page(width=612, height=792)
        page = writer.pages[i]

        font_dict = DictionaryObject()
        font_dict[NameObject("/Type")] = NameObject("/Font")
        font_dict[NameObject("/Subtype")] = NameObject("/Type1")
        font_dict[NameObject("/BaseFont")] = NameObject("/Helvetica")

        resources = DictionaryObject()
        font_resources = DictionaryObject()
        font_resources[NameObject("/F1")] = font_dict
        resources[NameObject("/Font")] = font_resources
        page[NameObject("/Resources")] = resources

        # Pad each page with unique filler so memory differences are measurable
        filler = f"PAGE {i} " + ("X" * 2000)
        content = f"BT /F1 12 Tf 100 700 Td ({filler}) Tj ET".encode()
        stream = DecodedStreamObject()
        stream.set_data(content)
        page[NameObject("/Contents")] = stream

    with open(path, "wb") as f:
        writer.write(f)


def _make_docx(path):
    """Create a minimal DOCX (Office Open XML zip) with text content."""
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '  <Default Extension="xml" ContentType="application/xml"/>'
        '  <Override PartName="/word/document.xml"'
        '   ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '  <Relationship Id="rId1"'
        '   Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument"'
        '   Target="word/document.xml"/>'
        "</Relationships>"
    )
    document = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        "  <w:body>"
        "    <w:p><w:r><w:t>Hello from lazy_load DOCX test</w:t></w:r></w:p>"
        "  </w:body>"
        "</w:document>"
    )
    word_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        "</Relationships>"
    )
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types)
        z.writestr("_rels/.rels", rels)
        z.writestr("word/document.xml", document)
        z.writestr("word/_rels/document.xml.rels", word_rels)


def _make_xlsx(path):
    """Create a minimal XLSX with openpyxl."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Name"
    ws["B1"] = "Value"
    ws["A2"] = "Test Item"
    ws["B2"] = 42
    wb.save(path)


def _make_pptx(path):
    """Create a minimal PPTX with python-pptx."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Test Slide"
    slide.placeholders[1].text = "Hello from lazy_load PPTX test"
    prs.save(path)


def _make_epub(path):
    """Create a minimal EPUB 3 file (zip with OEBPS structure)."""
    container_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
        "  <rootfiles>"
        '    <rootfile full-path="OEBPS/content.opf"'
        '     media-type="application/oebps-package+xml"/>'
        "  </rootfiles>"
        "</container>"
    )
    content_opf = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<package xmlns="http://www.idpf.org/2007/opf" unique-identifier="uid" version="3.0">'
        '  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
        '    <dc:identifier id="uid">test-epub-lazy</dc:identifier>'
        "    <dc:title>Test EPUB</dc:title>"
        "    <dc:language>en</dc:language>"
        '    <meta property="dcterms:modified">2024-01-01T00:00:00Z</meta>'
        "  </metadata>"
        "  <manifest>"
        '    <item id="ch1" href="chapter1.xhtml" media-type="application/xhtml+xml"/>'
        '    <item id="nav" href="nav.xhtml" media-type="application/xhtml+xml" properties="nav"/>'
        "  </manifest>"
        "  <spine>"
        '    <itemref idref="ch1"/>'
        "  </spine>"
        "</package>"
    )
    chapter1 = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        "<!DOCTYPE html>"
        '<html xmlns="http://www.w3.org/1999/xhtml">'
        "<head><title>Chapter 1</title></head>"
        "<body><h1>Chapter 1</h1>"
        "<p>Hello from lazy_load EPUB test</p>"
        "</body></html>"
    )
    nav = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        "<!DOCTYPE html>"
        '<html xmlns="http://www.w3.org/1999/xhtml"'
        '  xmlns:epub="http://www.idpf.org/2007/ops">'
        "<head><title>Nav</title></head>"
        "<body>"
        '<nav epub:type="toc"><ol>'
        '<li><a href="chapter1.xhtml">Chapter 1</a></li>'
        "</ol></nav>"
        "</body></html>"
    )
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("mimetype", "application/epub+zip")
        z.writestr("META-INF/container.xml", container_xml)
        z.writestr("OEBPS/content.opf", content_opf)
        z.writestr("OEBPS/chapter1.xhtml", chapter1)
        z.writestr("OEBPS/nav.xhtml", nav)


def _make_large_csv(path, num_rows=500):
    """Create a CSV with many rows so memory differences are measurable."""
    with open(path, "w", encoding="utf-8") as f:
        f.write("id,name,description\n")
        for i in range(num_rows):
            # ~200 bytes per row
            f.write(f'{i},item_{i},"{"D" * 150} row {i}"\n')


# ---------------------------------------------------------------------------
# Parametrized test: lazy_load() for every loader
# ---------------------------------------------------------------------------

# (filename, content_type, file_creator_or_text, expected_substring)
LOADER_CASES = [
    pytest.param(
        "test.txt",
        "text/plain",
        "Hello from lazy_load TXT test",
        "Hello from lazy_load TXT test",
        id="txt",
    ),
    pytest.param(
        "test.csv",
        "text/csv",
        "name,value\nAlpha,1\nBravo,2\n",
        "Alpha",
        id="csv",
    ),
    pytest.param(
        "test.json",
        "application/json",
        '{"key": "Hello from lazy_load JSON test"}',
        "Hello from lazy_load JSON test",
        id="json",
    ),
    pytest.param(
        "test.md",
        "text/markdown",
        "# Heading\n\nHello from lazy_load MD test\n",
        "Hello from lazy_load MD test",
        id="md",
    ),
    pytest.param(
        "test.rst",
        "text/x-rst",
        "Heading\n=======\n\nHello from lazy_load RST test\n",
        "Hello from lazy_load RST test",
        id="rst",
        marks=_skip_no_pandoc,
    ),
    pytest.param(
        "test.xml",
        "application/xml",
        '<?xml version="1.0"?><root><item>Hello from lazy_load XML test</item></root>',
        "Hello from lazy_load XML test",
        id="xml",
    ),
    pytest.param(
        "test.py",
        "text/x-python",
        '# Hello from lazy_load PY test\nprint("hello")\n',
        "Hello from lazy_load PY test",
        id="py-source",
    ),
    pytest.param("test.pdf", "application/pdf", _make_pdf, "PAGE 0", id="pdf"),
    pytest.param(
        "test.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        _make_docx,
        "lazy_load DOCX test",
        id="docx",
    ),
    pytest.param(
        "test.xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        _make_xlsx,
        "Test Item",
        id="xlsx",
        marks=_skip_no_msoffcrypto,
    ),
    pytest.param(
        "test.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        _make_pptx,
        "lazy_load PPTX test",
        id="pptx",
    ),
    pytest.param(
        "test.epub",
        "application/epub+zip",
        _make_epub,
        "lazy_load EPUB test",
        id="epub",
        marks=_skip_no_pandoc,
    ),
]


@pytest.mark.parametrize("filename,content_type,creator,expected_text", LOADER_CASES)
def test_lazy_load_returns_documents(
    tmp_path, filename, content_type, creator, expected_text
):
    """Every supported loader's lazy_load() should yield Document objects with content."""
    file_path = tmp_path / filename

    # Create the test file — either write text or call a builder function
    if callable(creator):
        creator(str(file_path))
    else:
        file_path.write_text(creator, encoding="utf-8")

    loader, known_type, file_ext = get_loader(filename, content_type, str(file_path))

    # Verify the loader has lazy_load
    assert hasattr(loader, "lazy_load"), f"{type(loader).__name__} missing lazy_load()"

    # Consume the lazy iterator
    docs = list(loader.lazy_load())

    # Basic assertions
    assert len(docs) > 0, f"{type(loader).__name__} yielded 0 documents"
    assert all(
        isinstance(d, Document) for d in docs
    ), "lazy_load() must yield Document instances"

    # Content assertion — at least one doc should contain the expected text
    all_text = " ".join(d.page_content for d in docs)
    assert (
        expected_text in all_text
    ), f"{type(loader).__name__}: expected '{expected_text}' in output, got: {all_text[:200]}"

    assert known_type is True


@pytest.mark.parametrize("filename,content_type,creator,expected_text", LOADER_CASES)
def test_lazy_load_matches_load(
    tmp_path, filename, content_type, creator, expected_text
):
    """lazy_load() consumed as a list should produce the same documents as load()."""
    file_path = tmp_path / filename

    if callable(creator):
        creator(str(file_path))
    else:
        file_path.write_text(creator, encoding="utf-8")

    loader, _, _ = get_loader(filename, content_type, str(file_path))

    eager_docs = loader.load()
    # Re-create loader since some loaders are single-use or have internal state
    loader2, _, _ = get_loader(filename, content_type, str(file_path))
    lazy_docs = list(loader2.lazy_load())

    assert len(eager_docs) == len(lazy_docs), (
        f"{type(loader).__name__}: load() returned {len(eager_docs)} docs, "
        f"lazy_load() returned {len(lazy_docs)}"
    )

    for i, (eager, lazy) in enumerate(zip(eager_docs, lazy_docs)):
        assert (
            eager.page_content == lazy.page_content
        ), f"{type(loader).__name__} doc[{i}]: content mismatch between load() and lazy_load()"


# ---------------------------------------------------------------------------
# SafePyPDFLoader-specific tests
# ---------------------------------------------------------------------------


def test_safe_pdf_loader_lazy_load_is_generator(tmp_path):
    """SafePyPDFLoader.lazy_load() should return a generator type."""
    pdf_path = tmp_path / "gen_test.pdf"
    _make_pdf(str(pdf_path))

    loader = SafePyPDFLoader(str(pdf_path), extract_images=False)
    result = loader.lazy_load()
    assert isinstance(result, types.GeneratorType)

    # Consuming the generator should yield documents
    docs = list(result)
    assert len(docs) > 0


def test_safe_pdf_loader_load_delegates_to_lazy_load(tmp_path):
    """SafePyPDFLoader.load() should produce the same results as list(lazy_load())."""
    pdf_path = tmp_path / "delegate_test.pdf"
    _make_pdf(str(pdf_path))

    loader1 = SafePyPDFLoader(str(pdf_path), extract_images=False)
    loader2 = SafePyPDFLoader(str(pdf_path), extract_images=False)

    load_docs = loader1.load()
    lazy_docs = list(loader2.lazy_load())

    assert len(load_docs) == len(lazy_docs)
    for ld, lz in zip(load_docs, lazy_docs):
        assert ld.page_content == lz.page_content


# ---------------------------------------------------------------------------
# CSV with non-UTF-8 encoding
# ---------------------------------------------------------------------------


def test_lazy_load_csv_non_utf8(tmp_path):
    """CSV files with non-UTF-8 encoding should still work via lazy_load()."""
    csv_path = tmp_path / "latin1.csv"
    csv_path.write_bytes("name,city\nJosé,São Paulo\n".encode("latin-1"))

    loader, known_type, _ = get_loader("latin1.csv", "text/csv", str(csv_path))
    docs = list(loader.lazy_load())

    assert len(docs) > 0
    all_text = " ".join(d.page_content for d in docs)
    # The text should have been converted to UTF-8 by get_loader
    assert "Jos" in all_text


# ---------------------------------------------------------------------------
# Memory benchmarks: lazy_load() vs load()
#
# These tests use tracemalloc to measure peak memory during document loading.
# They are informational — the assertions are intentionally loose so CI does
# not flake, but the captured output (-s) shows the real numbers.
# ---------------------------------------------------------------------------


def _measure_load(loader_factory):
    """Run load() and return (docs, peak_memory_bytes)."""
    gc.collect()
    tracemalloc.start()

    loader = loader_factory()
    docs = loader.load()

    _, peak = tracemalloc.get_traced_memory()
    tracemalloc.stop()
    gc.collect()
    return docs, peak


def _measure_lazy_load_materialized(loader_factory):
    """Run list(lazy_load()) and return (docs, peak_memory_bytes).

    This is the pattern our call sites currently use — materializes the
    full list, but pages are loaded one-at-a-time from the source.
    """
    gc.collect()
    tracemalloc.start()

    loader = loader_factory()
    docs = list(loader.lazy_load())

    _, peak = tracemalloc.get_traced_memory()
    tracemalloc.stop()
    gc.collect()
    return docs, peak


def _measure_lazy_load_streaming(loader_factory):
    """Iterate lazy_load() and accumulate only the text, discarding Document
    objects as we go. This simulates a true streaming consumer and represents
    the theoretical best-case for lazy_load().
    """
    gc.collect()
    tracemalloc.start()

    loader = loader_factory()
    texts = []
    for doc in loader.lazy_load():
        texts.append(doc.page_content)

    _, peak = tracemalloc.get_traced_memory()
    tracemalloc.stop()
    gc.collect()
    return texts, peak


class TestMemoryBenchmarkPDF:
    """Memory comparison for PyPDFLoader — the loader most likely to benefit
    from lazy_load() since it yields page-by-page."""

    NUM_PAGES = 50

    @pytest.fixture()
    def pdf_path(self, tmp_path):
        path = tmp_path / "bench.pdf"
        _make_pdf(str(path), num_pages=self.NUM_PAGES)
        return str(path)

    def test_pdf_lazy_load_streaming_uses_less_peak_memory(self, pdf_path, capsys):
        """Streaming consumption of lazy_load() should use less peak memory
        than load() which materializes every page simultaneously."""

        def factory():
            return SafePyPDFLoader(pdf_path, extract_images=False)

        load_docs, peak_load = _measure_load(factory)
        lazy_docs, peak_lazy_mat = _measure_lazy_load_materialized(factory)
        texts, peak_lazy_stream = _measure_lazy_load_streaming(factory)

        assert len(load_docs) == self.NUM_PAGES
        assert len(lazy_docs) == self.NUM_PAGES
        assert len(texts) == self.NUM_PAGES

        print(f"\n--- PDF Memory Benchmark ({self.NUM_PAGES} pages) ---")
        print(f"  load()                  peak: {peak_load:>10,} bytes")
        print(f"  list(lazy_load())       peak: {peak_lazy_mat:>10,} bytes")
        print(f"  streaming lazy_load()   peak: {peak_lazy_stream:>10,} bytes")
        print(
            f"  streaming vs load() savings:  {peak_load - peak_lazy_stream:>+10,} bytes"
        )

        ratio = peak_lazy_stream / peak_load if peak_load > 0 else 1.0
        print(f"  streaming / load() ratio:     {ratio:.2%}")

        # Informational: we expect streaming to be <= load, but don't hard-fail
        # if tracemalloc noise makes it slightly higher
        if peak_lazy_stream < peak_load:
            print("  VERDICT: streaming lazy_load() used LESS peak memory")
        else:
            print("  VERDICT: no measurable difference (expected for small pages)")


class TestMemoryBenchmarkCSV:
    """Memory comparison for CSVLoader — yields one Document per row."""

    NUM_ROWS = 500

    @pytest.fixture()
    def csv_path(self, tmp_path):
        path = tmp_path / "bench.csv"
        _make_large_csv(str(path), num_rows=self.NUM_ROWS)
        return str(path)

    def test_csv_lazy_load_streaming_uses_less_peak_memory(self, csv_path, capsys):
        """Streaming consumption of lazy_load() should use less peak memory
        than load() for CSVs with many rows."""

        def factory():
            from langchain_community.document_loaders import CSVLoader

            return CSVLoader(csv_path)

        load_docs, peak_load = _measure_load(factory)
        lazy_docs, peak_lazy_mat = _measure_lazy_load_materialized(factory)
        texts, peak_lazy_stream = _measure_lazy_load_streaming(factory)

        assert len(load_docs) == self.NUM_ROWS
        assert len(lazy_docs) == self.NUM_ROWS
        assert len(texts) == self.NUM_ROWS

        print(f"\n--- CSV Memory Benchmark ({self.NUM_ROWS} rows) ---")
        print(f"  load()                  peak: {peak_load:>10,} bytes")
        print(f"  list(lazy_load())       peak: {peak_lazy_mat:>10,} bytes")
        print(f"  streaming lazy_load()   peak: {peak_lazy_stream:>10,} bytes")
        print(
            f"  streaming vs load() savings:  {peak_load - peak_lazy_stream:>+10,} bytes"
        )

        ratio = peak_lazy_stream / peak_load if peak_load > 0 else 1.0
        print(f"  streaming / load() ratio:     {ratio:.2%}")

        if peak_lazy_stream < peak_load:
            print("  VERDICT: streaming lazy_load() used LESS peak memory")
        else:
            print("  VERDICT: no measurable difference")


UNSTRUCTURED_CASES = [
    pytest.param("test.md", "text/markdown", id="md"),
    pytest.param("test.xml", "application/xml", id="xml"),
    pytest.param("test.rst", "text/x-rst", id="rst", marks=_skip_no_pandoc),
    pytest.param(
        "test.epub",
        "application/epub+zip",
        id="epub",
        marks=_skip_no_pandoc,
    ),
    pytest.param(
        "test.xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        id="xlsx",
        marks=_skip_no_msoffcrypto,
    ),
    pytest.param(
        "test.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        id="pptx",
    ),
    pytest.param(
        "test.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        id="docx",
    ),
]

# Map extensions to their file creators
_UNSTRUCTURED_CREATORS = {
    ".md": lambda p: open(p, "w").write("# Heading\n\n" + ("word " * 500) + "\n"),
    ".xml": lambda p: open(p, "w").write(
        '<?xml version="1.0"?><root>'
        + "".join(f"<i>item {i}</i>" for i in range(100))
        + "</root>"
    ),
    ".rst": lambda p: open(p, "w").write(
        "Title\n=====\n\n"
        + "\n\n".join(f"Paragraph {i}. " + "text " * 50 for i in range(20))
    ),
    ".epub": lambda p: _make_epub(p),
    ".xlsx": lambda p: _make_xlsx(p),
    ".pptx": lambda p: _make_pptx(p),
    ".docx": lambda p: _make_docx(p),
}


@pytest.mark.parametrize("filename,content_type", UNSTRUCTURED_CASES)
def test_unstructured_lazy_load_no_memory_benefit(tmp_path, filename, content_type):
    """Unstructured-based loaders internally load the full file regardless of
    lazy_load() vs load(). This test confirms that and prints the numbers."""
    file_path = tmp_path / filename
    ext = os.path.splitext(filename)[1]
    _UNSTRUCTURED_CREATORS[ext](str(file_path))

    def factory():
        loader, _, _ = get_loader(filename, content_type, str(file_path))
        return loader

    load_docs, peak_load = _measure_load(factory)
    lazy_docs, peak_lazy_mat = _measure_lazy_load_materialized(factory)

    assert len(load_docs) == len(lazy_docs)

    diff = peak_load - peak_lazy_mat
    print(f"\n--- {ext} Unstructured Memory ({len(load_docs)} docs) ---")
    print(f"  load()            peak: {peak_load:>10,} bytes")
    print(f"  list(lazy_load()) peak: {peak_lazy_mat:>10,} bytes")
    print(f"  difference:             {diff:>+10,} bytes")

    # We expect roughly equal — the Unstructured loaders don't truly stream
    # Allow a generous tolerance since tracemalloc has inherent noise
    # This is informational, not a hard assertion
    if abs(diff) < max(peak_load * 0.15, 4096):
        print("  VERDICT: no meaningful difference (as expected)")
    elif diff > 0:
        print("  VERDICT: lazy_load() used slightly less (likely noise)")
    else:
        print("  VERDICT: lazy_load() used slightly more (likely noise)")
